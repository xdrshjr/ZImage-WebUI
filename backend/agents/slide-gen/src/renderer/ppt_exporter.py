"""
PPT export from HTML slides
"""

import logging
import re
import traceback
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PPTExporter:
    """Exports HTML slides to editable PowerPoint presentation"""
    
    def __init__(self):
        """Initialize PPT exporter"""
        # Aspect ratio dimensions in inches for PPT (standard presentation sizes)
        self.dimensions = {
            "16:9": {"width": Inches(10), "height": Inches(5.625)},
            "4:3": {"width": Inches(10), "height": Inches(7.5)},
            "16:10": {"width": Inches(10), "height": Inches(6.25)}
        }
        
        # Color scheme matching the HTML templates (Apple-like design)
        self.colors = {
            'title': RGBColor(29, 29, 31),  # #1d1d1f
            'text': RGBColor(29, 29, 31),   # #1d1d1f
            'section_header': RGBColor(134, 134, 139),  # #86868b
            'background': RGBColor(255, 255, 255),  # #ffffff
            'accent': RGBColor(0, 113, 227)  # #0071e3 - Apple blue
        }
    
    def _apply_color_scheme(self, scheme_colors: Dict[str, str]):
        """
        Apply color scheme by converting hex colors to RGBColor objects
        
        Args:
            scheme_colors: Dictionary with hex color values
        """
        def hex_to_rgb(hex_color: str) -> RGBColor:
            """Convert hex color string to RGBColor"""
            hex_color = hex_color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        
        # Update color mappings
        self.colors['background'] = hex_to_rgb(scheme_colors['background'])
        self.colors['text'] = hex_to_rgb(scheme_colors['text'])
        self.colors['title'] = hex_to_rgb(scheme_colors['text'])  # Title uses text color
        self.colors['accent'] = hex_to_rgb(scheme_colors['accent'])
        self.colors['section_header'] = hex_to_rgb(scheme_colors['header'])
        
        logger.debug(f"Color scheme applied: background={self.colors['background']}, "
                    f"text={self.colors['text']}, accent={self.colors['accent']}")
    
    def export_to_ppt(
        self,
        slides_data: List[Dict[str, Any]],
        output_path: Path,
        aspect_ratio: str = "16:9",
        color_scheme: str = "light_blue"
    ) -> bool:
        """
        Export slides data to PowerPoint presentation
        
        Args:
            slides_data: List of slide data dictionaries containing layout info
            output_path: Path for output PPTX file
            aspect_ratio: Aspect ratio for slide dimensions
            color_scheme: Color scheme name (e.g., "light_blue", "dark_slate")
            
        Returns:
            Success status
        """
        if not slides_data:
            logger.error("No slides data provided for PPT export")
            return False
        
        logger.info("=" * 60)
        logger.info("Starting PPT Generation")
        logger.info("=" * 60)
        logger.info(f"Total slides to export: {len(slides_data)}")
        logger.info(f"Aspect ratio: {aspect_ratio}")
        logger.info(f"Color scheme: {color_scheme}")
        logger.info(f"Output path: {output_path}")
        logger.debug(f"Output directory: {output_path.parent}")
        
        # Get color scheme colors from validators
        from src.utils.validators import ColorScheme
        scheme_colors = ColorScheme.get_scheme(color_scheme)
        logger.debug(f"Color scheme values: {scheme_colors}")
        
        # Convert hex colors to RGBColor objects
        self._apply_color_scheme(scheme_colors)
        logger.info(f"Applied {scheme_colors['name']} color scheme")
        logger.debug(f"  Background: {scheme_colors['background']}")
        logger.debug(f"  Text: {scheme_colors['text']}")
        logger.debug(f"  Accent: {scheme_colors['accent']}")
        
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide dimensions
            dims = self.dimensions.get(aspect_ratio, self.dimensions["16:9"])
            prs.slide_width = dims['width']
            prs.slide_height = dims['height']
            
            logger.debug(f"Presentation dimensions set: {prs.slide_width} x {prs.slide_height}")
            logger.info("Creating blank presentation template...")
            
            # Process each slide
            for i, slide_data in enumerate(slides_data):
                slide_number = slide_data.get('slide_number', i + 1)
                template_type = slide_data.get('template_type', 'title_and_content')
                
                logger.info(f"→ Processing Slide {slide_number}/{len(slides_data)}")
                logger.debug(f"  Template type: {template_type}")
                
                success = self._create_slide(prs, slide_data, aspect_ratio)
                
                if success:
                    logger.info(f"  ✓ Slide {slide_number} added successfully")
                else:
                    logger.warning(f"  ⚠ Slide {slide_number} added with issues")
            
            # Save presentation
            logger.info("Saving PowerPoint presentation...")
            logger.debug(f"Writing to file: {output_path}")
            
            prs.save(str(output_path))
            
            # Verify file was created
            if output_path.exists():
                file_size = output_path.stat().st_size
                logger.info("=" * 60)
                logger.info("✓ PPT Generation Completed Successfully")
                logger.info(f"  File: {output_path.name}")
                logger.info(f"  Size: {file_size / 1024:.2f} KB")
                logger.info(f"  Slides: {len(slides_data)}")
                logger.info(f"  Location: {output_path}")
                logger.info("=" * 60)
                return True
            else:
                logger.error("PPT file was not created")
                return False
                
        except ImportError as e:
            logger.error("python-pptx library is not installed")
            logger.error("Install it with: pip install python-pptx")
            logger.debug(f"Import error details: {str(e)}")
            return False
        except Exception as e:
            logger.error(f"Failed to create PowerPoint presentation: {str(e)}")
            import traceback
            logger.debug(f"Traceback: {traceback.format_exc()}")
            return False
    
    def _create_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        aspect_ratio: str
    ) -> bool:
        """
        Create a single slide in the presentation
        
        Args:
            prs: Presentation object
            slide_data: Slide data dictionary
            aspect_ratio: Aspect ratio
            
        Returns:
            Success status
        """
        try:
            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            slide_number = slide_data.get('slide_number', '?')
            template_type = slide_data.get('template_type', 'title_and_content')
            layout = slide_data.get('layout', {})
            
            logger.debug(f"  Creating slide with template: {template_type}")
            logger.debug(f"  Slide dimensions: {slide_width} x {slide_height}")
            
            # Add slide background (matching HTML template background color)
            logger.debug(f"  Adding slide background with color: {self.colors['background']}")
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['background']
            logger.debug(f"  ✓ Background color applied")
            
            # Extract layout data
            title = layout.get('title', '')
            content_blocks = layout.get('content_blocks', [])
            
            logger.debug(f"  Title: {title[:50]}..." if len(title) > 50 else f"  Title: {title}")
            logger.debug(f"  Content blocks: {len(content_blocks)}")
            
            # Add vertical bar decoration (matching HTML ::before pseudo-element)
            bar_left = Inches(0.5)
            bar_top = Inches(0.47)  # Optically aligned with title text
            bar_width = Pt(8)
            bar_height = Pt(33)  # Approximately 0.75em relative to 44pt font
            
            logger.debug(f"  Adding decorative vertical bar at left: {bar_left}, top: {bar_top}")
            bar_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left,
                bar_top,
                bar_width,
                bar_height
            )
            # Style the bar with accent color and rounded corners
            bar_fill = bar_shape.fill
            bar_fill.solid()
            bar_fill.fore_color.rgb = self.colors['accent']
            
            # Remove bar outline for clean appearance
            bar_line = bar_shape.line
            bar_line.fill.background()
            
            logger.debug(f"  ✓ Vertical bar added with accent color: {self.colors['accent']}")
            
            # Add title with Apple-style design
            title_left = Inches(0.5)
            title_top = Inches(0.4)
            title_width = slide_width - Inches(1.0)
            title_height = Inches(0.8)
            
            title_box = slide.shapes.add_textbox(
                title_left, title_top, title_width, title_height
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.margin_left = Pt(28)  # Space for decorative bar
            
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['title']
            title_para.alignment = PP_ALIGN.LEFT
            
            logger.debug(f"  ✓ Title added with {len(title)} characters")
            
            # Process content based on template type
            if template_type == 'title_and_content':
                self._add_title_and_content_layout(
                    slide, content_blocks, slide_width, slide_height
                )
            elif template_type == 'two_column':
                self._add_two_column_layout(
                    slide, content_blocks, slide_width, slide_height
                )
            elif template_type == 'image_focus':
                self._add_image_focus_layout(
                    slide, content_blocks, slide_width, slide_height
                )
            else:
                # Default to title_and_content
                logger.debug(f"  Unknown template type '{template_type}', using default")
                self._add_title_and_content_layout(
                    slide, content_blocks, slide_width, slide_height
                )
            
            return True
            
        except Exception as e:
            logger.error(f"  Failed to create slide {slide_data.get('slide_number', '?')}: {str(e)}")
            logger.debug(f"  Error details: {traceback.format_exc()}")
            return False
    
    def _add_title_and_content_layout(
        self,
        slide,
        content_blocks: List[Dict[str, Any]],
        slide_width,
        slide_height
    ):
        """Add content for title_and_content template"""
        logger.debug("    Applying title_and_content layout")
        
        content_top = Inches(1.5)
        content_left = Inches(0.8)
        content_width = slide_width - Inches(1.6)
        current_y = content_top
        
        # First, add any images (they go at the top)
        image_blocks = [b for b in content_blocks if b.get('type') == 'image_placeholder']
        text_blocks = [b for b in content_blocks if b.get('type') == 'text']
        
        logger.debug(f"    Found {len(image_blocks)} images and {len(text_blocks)} text blocks")
        
        for block in image_blocks:
            image_path = block.get('image_path')
            if image_path and Path(image_path).exists():
                try:
                    # Center the image
                    img_width = slide_width * 0.6
                    img_height = Inches(2.0)
                    img_left = (slide_width - img_width) / 2
                    
                    slide.shapes.add_picture(
                        str(image_path),
                        img_left,
                        current_y,
                        width=img_width,
                        height=img_height
                    )
                    current_y += img_height + Inches(0.3)
                    logger.debug(f"    ✓ Image added: {Path(image_path).name}")
                except Exception as e:
                    logger.warning(f"    ⚠ Failed to add image: {str(e)}")
        
        # Then add text blocks
        for i, block in enumerate(text_blocks):
            section_title = block.get('section_title', '')
            content = block.get('content', '')
            
            if not content:
                continue
            
            # Calculate available space
            available_height = slide_height - current_y - Inches(0.5)
            
            if available_height < Inches(0.5):
                logger.debug(f"    ⚠ Insufficient space for text block {i+1}, skipping")
                break
            
            # Add text box
            text_box_height = min(Inches(1.5), available_height)
            text_box = slide.shapes.add_textbox(
                content_left,
                current_y,
                content_width,
                text_box_height
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            # Add section title if present
            if section_title:
                para = text_frame.paragraphs[0]
                para.text = section_title
                para.font.size = Pt(20)
                para.font.bold = True
                para.font.color.rgb = self.colors['section_header']
                para.alignment = PP_ALIGN.LEFT
                para.space_after = Pt(8)
                
                # Add content in new paragraph
                para = text_frame.add_paragraph()
            else:
                para = text_frame.paragraphs[0]
            
            # Clean and format content
            cleaned_content = self._clean_text_content(content)
            para.text = cleaned_content
            para.font.size = Pt(16)
            para.font.color.rgb = self.colors['text']
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 1.3
            
            current_y += text_box_height + Inches(0.2)
            logger.debug(f"    ✓ Text block {i+1} added")
    
    def _add_two_column_layout(
        self,
        slide,
        content_blocks: List[Dict[str, Any]],
        slide_width,
        slide_height
    ):
        """Add content for two_column template"""
        logger.debug("    Applying two_column layout")
        
        # Divide content into left and right columns based on position.x
        left_blocks = []
        right_blocks = []
        
        midpoint = 1920 / 2  # Based on HTML width
        
        for block in content_blocks:
            position = block.get('position', {})
            x_pos = position.get('x', 0)
            
            if x_pos < midpoint:
                left_blocks.append(block)
            else:
                right_blocks.append(block)
        
        logger.debug(f"    Left column: {len(left_blocks)} blocks, Right column: {len(right_blocks)} blocks")
        
        # Define column dimensions
        column_width = (slide_width - Inches(2.0)) / 2
        content_top = Inches(1.5)
        left_column_left = Inches(0.6)
        right_column_left = left_column_left + column_width + Inches(0.4)
        
        # Process left column
        self._add_column_content(
            slide, left_blocks, left_column_left, content_top,
            column_width, slide_height, "left"
        )
        
        # Process right column
        self._add_column_content(
            slide, right_blocks, right_column_left, content_top,
            column_width, slide_height, "right"
        )
    
    def _add_column_content(
        self,
        slide,
        blocks: List[Dict[str, Any]],
        column_left,
        content_top,
        column_width,
        slide_height,
        column_name: str
    ):
        """Add content to a column"""
        current_y = content_top
        
        for i, block in enumerate(blocks):
            block_type = block.get('type')
            
            if block_type == 'image_placeholder':
                image_path = block.get('image_path')
                if image_path and Path(image_path).exists():
                    try:
                        img_height = Inches(1.8)
                        
                        if current_y + img_height > slide_height - Inches(0.5):
                            logger.debug(f"    ⚠ Insufficient space for image in {column_name} column")
                            break
                        
                        slide.shapes.add_picture(
                            str(image_path),
                            column_left,
                            current_y,
                            width=column_width,
                            height=img_height
                        )
                        current_y += img_height + Inches(0.2)
                        logger.debug(f"    ✓ Image added to {column_name} column")
                    except Exception as e:
                        logger.warning(f"    ⚠ Failed to add image to {column_name} column: {str(e)}")
            
            elif block_type == 'text':
                section_title = block.get('section_title', '')
                content = block.get('content', '')
                
                if not content:
                    continue
                
                available_height = slide_height - current_y - Inches(0.5)
                if available_height < Inches(0.3):
                    logger.debug(f"    ⚠ Insufficient space for text in {column_name} column")
                    break
                
                text_box_height = min(Inches(1.2), available_height)
                text_box = slide.shapes.add_textbox(
                    column_left,
                    current_y,
                    column_width,
                    text_box_height
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                
                if section_title:
                    para = text_frame.paragraphs[0]
                    para.text = section_title
                    para.font.size = Pt(18)
                    para.font.bold = True
                    para.font.color.rgb = self.colors['section_header']
                    para.alignment = PP_ALIGN.LEFT
                    para.space_after = Pt(6)
                    
                    para = text_frame.add_paragraph()
                else:
                    para = text_frame.paragraphs[0]
                
                cleaned_content = self._clean_text_content(content)
                para.text = cleaned_content
                para.font.size = Pt(14)
                para.font.color.rgb = self.colors['text']
                para.alignment = PP_ALIGN.LEFT
                para.line_spacing = 1.3
                
                current_y += text_box_height + Inches(0.15)
                logger.debug(f"    ✓ Text block added to {column_name} column")
    
    def _add_image_focus_layout(
        self,
        slide,
        content_blocks: List[Dict[str, Any]],
        slide_width,
        slide_height
    ):
        """Add content for image_focus template (absolute positioning)"""
        logger.debug("    Applying image_focus layout")
        
        # HTML dimensions for conversion
        html_width = 1920
        html_height = 1080
        
        for i, block in enumerate(content_blocks):
            block_type = block.get('type')
            position = block.get('position', {})
            
            # Convert HTML pixel coordinates to PPT inches
            x_ratio = position.get('x', 0) / html_width
            y_ratio = position.get('y', 0) / html_height
            width_ratio = position.get('width', 200) / html_width
            height_ratio = position.get('height', 200) / html_height
            
            ppt_left = slide_width * x_ratio
            ppt_top = slide_height * y_ratio
            ppt_width = slide_width * width_ratio
            ppt_height = slide_height * height_ratio
            
            if block_type == 'image_placeholder':
                image_path = block.get('image_path')
                if image_path and Path(image_path).exists():
                    try:
                        slide.shapes.add_picture(
                            str(image_path),
                            ppt_left,
                            ppt_top,
                            width=ppt_width,
                            height=ppt_height
                        )
                        logger.debug(f"    ✓ Image positioned at ({x_ratio:.2%}, {y_ratio:.2%})")
                    except Exception as e:
                        logger.warning(f"    ⚠ Failed to add positioned image: {str(e)}")
            
            elif block_type == 'text':
                section_title = block.get('section_title', '')
                content = block.get('content', '')
                
                if not content:
                    continue
                
                text_box = slide.shapes.add_textbox(
                    ppt_left,
                    ppt_top,
                    ppt_width,
                    ppt_height
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                
                if section_title:
                    para = text_frame.paragraphs[0]
                    para.text = section_title
                    para.font.size = Pt(16)
                    para.font.bold = True
                    para.font.color.rgb = self.colors['section_header']
                    para.alignment = PP_ALIGN.LEFT
                    para.space_after = Pt(6)
                    
                    para = text_frame.add_paragraph()
                else:
                    para = text_frame.paragraphs[0]
                
                cleaned_content = self._clean_text_content(content)
                para.text = cleaned_content
                para.font.size = Pt(14)
                para.font.color.rgb = self.colors['text']
                para.alignment = PP_ALIGN.LEFT
                para.line_spacing = 1.2
                
                logger.debug(f"    ✓ Text positioned at ({x_ratio:.2%}, {y_ratio:.2%})")
    
    def _clean_text_content(self, text: str) -> str:
        """
        Clean text content from HTML formatting
        
        Args:
            text: Raw text content
            
        Returns:
            Cleaned text
        """
        if not text:
            return ""
        
        # Remove excessive whitespace and newlines
        text = re.sub(r'\n\s*\n', '\n', text)
        text = text.strip()
        
        return text

